import os
import re
import ffmpeg
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches
import pyttsx3
import tempfile

# ============ 路径配置 =============
os.environ['PATH'] += os.pathsep + r'D:\ffmpeg-7.1.1-full_build\bin'

# ============ PPT生成函数 =============
def create_vertical_ppt(image_folder, output_ppt):
    """生成竖版PPT（仅图片）"""
    prs = Presentation()
    prs.slide_width = Inches(9)
    prs.slide_height = Inches(16)

    valid_extensions = ['.jpg', '.jpeg', '.png', '.bmp']
    image_paths = sorted([
        os.path.join(image_folder, f) 
        for f in os.listdir(image_folder) 
        if os.path.splitext(f)[1].lower() in valid_extensions
    ])

    for img_path in image_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), 
                               width=prs.slide_width, 
                               height=prs.slide_height)

    prs.save(output_ppt)
    print(f"PPT已生成：{os.path.abspath(output_ppt)}")

# ============ 从HTML提取文本 =============
def extract_text_from_html(html_path):
    """从HTML文件的<div class="content">中提取<p>标签文本"""
    with open(html_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')
    
    content_div = soup.find('div', class_='content')
    if not content_div:
        return ""
    
    paragraphs = content_div.find_all('p')
    text = ' '.join(p.get_text(strip=True) for p in paragraphs if p.get_text(strip=True))
    
    # 简单清洗
    text = re.sub(r'\s+', ' ', text)
    text = text.replace('&nbsp;', ' ')
    return text
def text_to_speech(text, audio_path):
    """将文本转换为语音"""
    engine = pyttsx3.init()
    voices = engine.getProperty('voices')
    for voice in voices:
        if "Chinese" in voice.name:
            engine.setProperty('voice', voice.id)
            break
    engine.save_to_file(text, audio_path)
    engine.runAndWait()

def generate_video_with_effects(image_paths, audio_paths, output_video):
    """生成带有效果的视频"""
    temp_dir = tempfile.mkdtemp()
    video_segments = []

    target_width = 1080
    target_height = 1920

    image_paths = sorted(image_paths)
    audio_paths = sorted(audio_paths)

    audio_durations = []
    for audio_path in audio_paths:
        duration = float(ffmpeg.probe(audio_path)['format']['duration'])
        audio_durations.append(duration)

    for i, (img_path, audio_path, duration) in enumerate(zip(image_paths, audio_paths, audio_durations)):
        output_segment = os.path.join(temp_dir, f"segment_{i}.mp4")

        (
            ffmpeg
            .input(img_path, loop=1, t=duration)
            .filter('scale', width=target_width, height=target_height, force_original_aspect_ratio='decrease')
            .filter('pad', width=target_width, height=target_height, x='(ow-iw)/2', y='(oh-ih)/2', color='black')
            .output(
                ffmpeg.input(audio_path),
                output_segment,
                vcodec='libx264',
                acodec='aac',
                pix_fmt='yuv420p',
                **{'aspect': '9:16'},
                y=None
            )
            .run(overwrite_output=True)
        )
        video_segments.append(output_segment)

    concat_file = os.path.join(temp_dir, "concat_list.txt")
    with open(concat_file, 'w') as f:
        for segment in video_segments:
            f.write(f"file '{segment}'\n")

    (
        ffmpeg
        .input(concat_file, format='concat', safe=0)
        .output(output_video, c='copy', movflags='+faststart')
        .run(overwrite_output=True)
    )
# ============ 主流程 ============
def process_posters(html_folder, image_folder, output_ppt=None, output_video=None):
    """处理分开存放的HTML和图片文件"""
    if output_ppt is None:
        output_ppt = os.path.join(image_folder, "竖版海报.pptx")
    if output_video is None:
        output_video = os.path.join(image_folder, "output.mp4")

    # 1. 生成PPT（仅图片）
    create_vertical_ppt(image_folder, output_ppt)

    # 2. 获取HTML和图片文件列表（按文件名排序）
    html_files = sorted([
        os.path.join(html_folder, f) 
        for f in os.listdir(html_folder) 
        if f.lower().endswith('.html')
    ])

    image_extensions = ['.jpg', '.jpeg', '.png', '.bmp']
    image_paths = sorted([
        os.path.join(image_folder, f) 
        for f in os.listdir(image_folder)
        if os.path.splitext(f)[1].lower() in image_extensions
    ])

    # 3. 检查文件数量是否匹配
    if len(html_files) != len(image_paths):
        raise ValueError(
            f"HTML文件数量({len(html_files)})和图片数量({len(image_paths)})不匹配！\n"
            f"HTML文件: {html_files}\n图片文件: {image_paths}"
        )

    # 4. 处理每个HTML+图片对
    with tempfile.TemporaryDirectory() as tmpdir:
        audio_paths = []
        for i, (html_file, img_path) in enumerate(zip(html_files, image_paths)):
            # 提取文本并生成语音
            text = extract_text_from_html(html_file)
            audio_path = os.path.join(tmpdir, f"audio_{i}.mp3")
            text_to_speech(text, audio_path)
            audio_paths.append(audio_path)

        # 5. 生成视频
        generate_video_with_effects(image_paths, audio_paths, output_video)

    print(f"✅ 视频已生成：{output_video}")

# 使用示例
# if __name__ == "__main__":
#     html_folder = r"D:\代码\html"   # HTML文件夹路径
#     image_folder = r"D:\代码\posters"     # 图片文件夹路径
#     output_ppt = r"D:\代码\pptx\海报.pptx"     # PPT保存位置
#     output_video = r"D:\代码\audio\final_video.mp4" 
#     process_posters(html_folder, image_folder, output_ppt, output_video)